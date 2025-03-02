import os
import io
import json
import csv
import logging
import numpy as np
import spacy
import tiktoken
import fitz  # PyMuPDF for PDF processing
from langdetect import detect, DetectorFactory
from sklearn.cluster import KMeans
from typing import List, Dict, Tuple, Optional

# For OCR-based image text extraction.
from PIL import Image
import pytesseract

# Configure logging to display info and error messages.
logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

# Fix randomness in language detection.
DetectorFactory.seed = 0


class DocumentTextEmbeddingsExtractor:
    """
    A class to extract text from various document types (PDF, image, DOC/DOCX, PPT/PPTX, CSV, XLSX),
    split the text into embedding-friendly chunks, and optionally write the results to a JSON file.

    For PDFs, the extraction process combines both native text extraction and OCR text from embedded images.
    """

    def __init__(
        self, model: str = "openai", sliding_window_sentences: int = 2
    ) -> None:
        """
        Initializes the extractor with model-specific settings.

        Args:
            model (str): Target embedding model identifier (e.g., "openai", "bert", "cohere").
            sliding_window_sentences (int): Number of overlapping sentences to preserve for context.
        """
        self.model: str = model
        self.sliding_window_sentences: int = sliding_window_sentences

        # Model-specific chunk sizes (token limits)
        self.chunk_sizes: Dict[str, int] = {
            "openai": 8192,  # OpenAI's embedding model token limit
            "bert": 512,  # BERT models (e.g., FAISS, Pinecone)
            "cohere": 1024,  # Cohere's embeddings
        }

        # Load spaCy NLP models.
        self.nlp_models: Dict[str, spacy.language.Language] = {
            "en": spacy.load(
                "en_core_web_trf"
            ),  # High-accuracy transformer model for English
            "xx": spacy.load("xx_sent_ud_sm"),  # Multilingual sentence segmentation
        }

        # Tokenizer selection for different models.
        self.tokenizers: Dict[str, any] = {
            "openai": tiktoken.encoding_for_model("text-embedding-ada-002")
        }

    def _detect_language(self, text: str) -> str:
        """
        Detects the language of a given text.

        Args:
            text (str): The text for language detection.

        Returns:
            str: Detected language code, or "xx" on failure.
        """
        try:
            return detect(text)
        except Exception as e:
            logging.warning("Language detection failed (%s). Falling back to 'xx'.", e)
            return "xx"

    def _detect_columns(
        self, text_blocks: List[Dict], num_columns: Optional[int] = None
    ) -> Tuple[List[Dict], int]:
        """
        Automatically detects the number of columns on a PDF page using K-Means clustering.

        Args:
            text_blocks (List[Dict]): List of text block dictionaries extracted from a PDF page.
            num_columns (Optional[int]): Manual override for number of columns.

        Returns:
            Tuple[List[Dict], int]: The (possibly re-ordered) text blocks with a "column" key added,
            and the number of detected columns.
        """
        if not text_blocks:
            return text_blocks, 0

        # Extract x-coordinates from blocks that have a bounding box.
        x_coords = np.array(
            [[block["bbox"][0]] for block in text_blocks if "bbox" in block]
        )
        if x_coords.size == 0:
            return text_blocks, 1

        if num_columns is None:
            num_columns = min(4, len(np.unique(x_coords)))

        kmeans = KMeans(n_clusters=num_columns, random_state=0, n_init=10)
        kmeans.fit(x_coords)

        for block, label in zip(text_blocks, kmeans.labels_):
            block["column"] = label

        text_blocks.sort(key=lambda b: (b.get("column", 0), b.get("bbox", [0, 0])[1]))
        return text_blocks, num_columns

    def _count_tokens(self, text: str, model: Optional[str] = None) -> int:
        """
        Counts tokens in the given text using the specified model's tokenizer.

        Args:
            text (str): Text to tokenize.
            model (Optional[str]): Embedding model identifier.

        Returns:
            int: The number of tokens.
        """
        if model is None:
            model = self.model
        tokenizer = self.tokenizers.get(model, self.tokenizers["openai"])
        return len(tokenizer.encode(text))

    def _split_text_for_embeddings(
        self, text: str, lang: str, model: Optional[str] = None
    ) -> List[str]:
        """
        Splits text into chunks that do not exceed the embedding model's token limit.
        It preserves whole sentences and retains a sliding window overlap between chunks.

        Args:
            text (str): The full text to split.
            lang (str): Language code to select the appropriate NLP model.
            model (Optional[str]): Target embedding model.

        Returns:
            List[str]: List of text chunks.
        """
        if model is None:
            model = self.model
        if not text.strip():
            return []

        nlp = self.nlp_models.get(lang, self.nlp_models["xx"])
        doc = nlp(text)
        chunk_size = self.chunk_sizes.get(model, 512)
        chunks: List[str] = []
        current_chunk: List[str] = []
        current_length = 0

        sentences: List[str] = [
            sent.text.strip() for sent in doc.sents if sent.text.strip()
        ]

        for sent_text in sentences:
            sent_length = self._count_tokens(sent_text, model)
            if current_length + sent_length > chunk_size:
                chunks.append(" ".join(current_chunk))
                current_chunk = current_chunk[-self.sliding_window_sentences :]
                current_length = self._count_tokens(" ".join(current_chunk), model)
            current_chunk.append(sent_text)
            current_length += sent_length

        if current_chunk:
            chunks.append(" ".join(current_chunk))
        return chunks

    def _extract_text_from_blocks(self, text_blocks: List[Dict]) -> str:
        """
        Extracts and concatenates text from a list of PDF text blocks.

        Args:
            text_blocks (List[Dict]): List of text block dictionaries.

        Returns:
            str: Combined text from the blocks.
        """
        texts = []
        for block in text_blocks:
            if "lines" in block:
                for line in block.get("lines", []):
                    line_text = " ".join(
                        span["text"]
                        for span in line.get("spans", [])
                        if span.get("text", "").strip()
                    )
                    if line_text:
                        texts.append(line_text)
        return " ".join(texts).strip()

    def _extract_text_by_words(self, page: fitz.Page, tolerance: float = 3.0) -> str:
        """
        Extracts text from the page using a word-level strategy. It groups individual words
        into lines based on their vertical proximity. This method can better reconstruct the
        original layout in some cases.

        Args:
            page (fitz.Page): The PDF page to extract text from.
            tolerance (float): Maximum difference in the y-coordinate to consider words on the same line.

        Returns:
            str: The reconstructed text with lines separated by newlines.
        """
        words = page.get_text(
            "words"
        )  # Each word: (x0, y0, x1, y1, "word", block_no, line_no, word_no)
        if not words:
            return ""

        words.sort(key=lambda w: (w[1], w[0]))
        lines = []
        current_line = []
        last_y = None

        for w in words:
            if last_y is None or abs(w[1] - last_y) < tolerance:
                current_line.append(w)
            else:
                current_line.sort(key=lambda x: x[0])
                line_text = " ".join(word[4] for word in current_line)
                lines.append(line_text)
                current_line = [w]
            last_y = w[1]

        if current_line:
            current_line.sort(key=lambda x: x[0])
            line_text = " ".join(word[4] for word in current_line)
            lines.append(line_text)

        return "\n".join(lines)

    def _extract_text_from_images(self, page: fitz.Page) -> str:
        """
        Extracts text from any images embedded in the PDF page using OCR.

        Args:
            page (fitz.Page): The PDF page to process.

        Returns:
            str: The concatenated OCR text from all images on the page.
        """
        texts = []
        images = page.get_images(full=True)
        if not images:
            return ""
        doc = page.parent
        for img in images:
            xref = img[0]
            try:
                image_dict = doc.extract_image(xref)
                image_bytes = image_dict["image"]
                image = Image.open(io.BytesIO(image_bytes))
                ocr_text = pytesseract.image_to_string(image)
                if ocr_text.strip():
                    texts.append(ocr_text.strip())
            except Exception as e:
                logging.warning("Failed to extract or OCR image (xref %s): %s", xref, e)
        return "\n".join(texts)

    def _extract_text_from_page(self, page: fitz.Page) -> str:
        """
        Attempts to reconstruct the original text from a PDF page using two methods:
        one based on text blocks and another based on individual words.

        Args:
            page (fitz.Page): The PDF page to extract text from.

        Returns:
            str: The reconstructed text from the page.
        """
        try:
            page_dict = page.get_text("dict")
            blocks = page_dict.get("blocks", [])
            blocks, _ = self._detect_columns(blocks)
            text_by_blocks = self._extract_text_from_blocks(blocks)
        except Exception as e:
            logging.warning("Block extraction failed: %s", e)
            text_by_blocks = ""

        try:
            text_by_words = self._extract_text_by_words(page)
        except Exception as e:
            logging.warning("Word-level extraction failed: %s", e)
            text_by_words = ""

        return (
            text_by_words
            if len(text_by_words.split()) >= len(text_by_blocks.split())
            else text_by_blocks
        )

    def _process_page(self, page: fitz.Page, page_num: int) -> List[Dict]:
        """
        Processes a single PDF page: extracts text (including OCR from images), detects language,
        splits text into chunks, and returns a list of dictionaries for each chunk.

        Args:
            page (fitz.Page): The PDF page to process.
            page_num (int): The page number.

        Returns:
            List[Dict]: List of dictionaries containing chunk data.
        """
        text_from_page = self._extract_text_from_page(page)
        text_from_images = self._extract_text_from_images(page)
        full_text = (
            (text_from_page + "\n[Image Content]\n" + text_from_images)
            if text_from_images
            else text_from_page
        )

        if not full_text.strip():
            logging.debug("Page %s is empty or contains no extractable text.", page_num)
            return []

        detected_lang = self._detect_language(full_text[:500])
        chunks = self._split_text_for_embeddings(full_text, detected_lang)
        chunks_data = []
        for i, chunk in enumerate(chunks):
            chunks_data.append(
                {
                    "page": page_num,
                    "chunk_id": f"{page_num}-{i}",
                    "language": detected_lang,
                    "tokens": self._count_tokens(chunk),
                    "text": chunk,
                }
            )
        return chunks_data

    # ==============================
    # PDF & Image Extraction Methods
    # ==============================

    def extract_text_from_pdf(self, pdf_path: str, output_json_path: str) -> Dict:
        """
        Extracts text from each page of a PDF (including OCR text for images), splits it into chunks,
        and writes the result to a JSON file.

        Args:
            pdf_path (str): Path to the PDF file.
            output_json_path (str): Path where the JSON output will be saved.

        Returns:
            Dict: A dictionary containing all the text chunks, model used, and source type.
        """
        data = {"chunks": [], "model": self.model, "source": "pdf"}
        try:
            with fitz.open(pdf_path) as doc:
                for page_num, page in enumerate(doc, start=1):
                    logging.info("Processing PDF page %s", page_num)
                    page_chunks = self._process_page(page, page_num)
                    data["chunks"].extend(page_chunks)
        except Exception as e:
            logging.error("Failed to process PDF file '%s': %s", pdf_path, e)
            raise e

        try:
            with open(output_json_path, "w", encoding="utf-8") as json_file:
                json.dump(data, json_file, indent=4, ensure_ascii=False)
            logging.info("Optimized text chunks saved to %s", output_json_path)
        except Exception as e:
            logging.error(
                "Failed to write JSON output to '%s': %s", output_json_path, e
            )
            raise e

        return data

    def extract_text_from_image(
        self, image_path: str, output_json_path: Optional[str] = None
    ) -> Dict:
        """
        Extracts text from an image file using OCR, splits it into chunks,
        and optionally writes the result to a JSON file.

        Args:
            image_path (str): Path to the image file.
            output_json_path (Optional[str]): Path to save the JSON output.

        Returns:
            Dict: A dictionary containing the extracted text chunks and model info.
        """
        try:
            image = Image.open(image_path)
        except Exception as e:
            logging.error("Failed to open image '%s': %s", image_path, e)
            raise e

        text = pytesseract.image_to_string(image)
        if not text.strip():
            logging.warning("No text found in image '%s'.", image_path)
            return {}

        detected_lang = self._detect_language(text[:500])
        chunks = self._split_text_for_embeddings(text, detected_lang)
        chunks_data = []
        for i, chunk in enumerate(chunks):
            chunks_data.append(
                {
                    "page": 1,
                    "chunk_id": f"1-{i}",
                    "language": detected_lang,
                    "tokens": self._count_tokens(chunk),
                    "text": chunk,
                }
            )
        data = {"chunks": chunks_data, "model": self.model, "source": "image"}

        if output_json_path:
            try:
                with open(output_json_path, "w", encoding="utf-8") as json_file:
                    json.dump(data, json_file, indent=4, ensure_ascii=False)
                logging.info("Extracted text from image saved to %s", output_json_path)
            except Exception as e:
                logging.error(
                    "Failed to write JSON output to '%s': %s", output_json_path, e
                )
                raise e
        return data

    # ==============================
    # DOC / DOCX Extraction Methods
    # ==============================

    def extract_text_from_docx(
        self, file_path: str, output_json_path: Optional[str] = None
    ) -> Dict:
        """
        Extracts text from a DOCX file using python-docx, splits it into chunks,
        and optionally writes the result to a JSON file.

        Args:
            file_path (str): Path to the DOCX file.
            output_json_path (Optional[str]): Path to save the JSON output.

        Returns:
            Dict: A dictionary containing the extracted text chunks and model info.
        """
        try:
            import docx
        except ImportError:
            logging.error(
                "python-docx is required for DOCX extraction. Please install it."
            )
            raise

        document = docx.Document(file_path)
        full_text = "\n".join(
            [para.text for para in document.paragraphs if para.text.strip()]
        )
        detected_lang = self._detect_language(full_text[:500])
        chunks = self._split_text_for_embeddings(full_text, detected_lang)
        chunks_data = []
        for i, chunk in enumerate(chunks):
            chunks_data.append(
                {
                    "page": 1,
                    "chunk_id": f"1-{i}",
                    "language": detected_lang,
                    "tokens": self._count_tokens(chunk),
                    "text": chunk,
                }
            )
        data = {"chunks": chunks_data, "model": self.model, "source": "docx"}

        if output_json_path:
            try:
                with open(output_json_path, "w", encoding="utf-8") as json_file:
                    json.dump(data, json_file, indent=4, ensure_ascii=False)
                logging.info("Extracted text from DOCX saved to %s", output_json_path)
            except Exception as e:
                logging.error(
                    "Failed to write JSON output to '%s': %s", output_json_path, e
                )
                raise e
        return data

    # ==============================
    # PPT / PPTX Extraction Methods
    # ==============================

    def extract_text_from_pptx(
        self, file_path: str, output_json_path: Optional[str] = None
    ) -> Dict:
        """
        Extracts text from a PPTX file using python-pptx, splits it into chunks per slide,
        and optionally writes the result to a JSON file.

        Args:
            file_path (str): Path to the PPTX file.
            output_json_path (Optional[str]): Path to save the JSON output.

        Returns:
            Dict: A dictionary containing the extracted text chunks and model info.
        """
        try:
            from pptx import Presentation
        except ImportError:
            logging.error(
                "python-pptx is required for PPTX extraction. Please install it."
            )
            raise

        prs = Presentation(file_path)
        chunks_data = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text)
            full_text = "\n".join(slide_texts)
            detected_lang = self._detect_language(full_text[:500])
            chunks = self._split_text_for_embeddings(full_text, detected_lang)
            for j, chunk in enumerate(chunks):
                chunks_data.append(
                    {
                        "page": slide_num,
                        "chunk_id": f"{slide_num}-{j}",
                        "language": detected_lang,
                        "tokens": self._count_tokens(chunk),
                        "text": chunk,
                    }
                )
        data = {"chunks": chunks_data, "model": self.model, "source": "pptx"}

        if output_json_path:
            try:
                with open(output_json_path, "w", encoding="utf-8") as json_file:
                    json.dump(data, json_file, indent=4, ensure_ascii=False)
                logging.info("Extracted text from PPTX saved to %s", output_json_path)
            except Exception as e:
                logging.error(
                    "Failed to write JSON output to '%s': %s", output_json_path, e
                )
                raise e
        return data

    # ==============================
    # CSV and XLSX Extraction Methods
    # ==============================

    def extract_text_from_csv(
        self, file_path: str, output_json_path: Optional[str] = None
    ) -> Dict:
        """
        Extracts text from a CSV file by concatenating all rows, splits it into chunks,
        and optionally writes the result to a JSON file.

        Args:
            file_path (str): Path to the CSV file.
            output_json_path (Optional[str]): Path to save the JSON output.

        Returns:
            Dict: A dictionary containing the extracted text chunks and model info.
        """
        rows = []
        try:
            with open(file_path, newline="", encoding="utf-8") as csvfile:
                reader = csv.reader(csvfile)
                for row in reader:
                    rows.append(", ".join(row))
        except Exception as e:
            logging.error("Failed to read CSV file '%s': %s", file_path, e)
            raise e

        full_text = "\n".join(rows)
        detected_lang = self._detect_language(full_text[:500])
        chunks = self._split_text_for_embeddings(full_text, detected_lang)
        chunks_data = []
        for i, chunk in enumerate(chunks):
            chunks_data.append(
                {
                    "page": 1,
                    "chunk_id": f"1-{i}",
                    "language": detected_lang,
                    "tokens": self._count_tokens(chunk),
                    "text": chunk,
                }
            )
        data = {"chunks": chunks_data, "model": self.model, "source": "csv"}

        if output_json_path:
            try:
                with open(output_json_path, "w", encoding="utf-8") as json_file:
                    json.dump(data, json_file, indent=4, ensure_ascii=False)
                logging.info("Extracted text from CSV saved to %s", output_json_path)
            except Exception as e:
                logging.error(
                    "Failed to write JSON output to '%s': %s", output_json_path, e
                )
                raise e
        return data

    def extract_text_from_xlsx(
        self, file_path: str, output_json_path: Optional[str] = None
    ) -> Dict:
        """
        Extracts text from an XLSX file by reading each sheet and concatenating cell values,
        splits it into chunks (each sheet is treated as a separate 'page'),
        and optionally writes the result to a JSON file.

        Args:
            file_path (str): Path to the XLSX file.
            output_json_path (Optional[str]): Path to save the JSON output.

        Returns:
            Dict: A dictionary containing the extracted text chunks and model info.
        """
        try:
            from openpyxl import load_workbook
        except ImportError:
            logging.error(
                "openpyxl is required for XLSX extraction. Please install it."
            )
            raise

        try:
            wb = load_workbook(file_path, read_only=True)
        except Exception as e:
            logging.error("Failed to open XLSX file '%s': %s", file_path, e)
            raise e

        chunks_data = []
        page_num = 1
        for sheet in wb.worksheets:
            sheet_texts = []
            for row in sheet.iter_rows(values_only=True):
                row_text = ", ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    sheet_texts.append(row_text)
            full_text = "\n".join(sheet_texts)
            detected_lang = self._detect_language(full_text[:500])
            chunks = self._split_text_for_embeddings(full_text, detected_lang)
            for j, chunk in enumerate(chunks):
                chunks_data.append(
                    {
                        "page": page_num,
                        "chunk_id": f"{page_num}-{j}",
                        "language": detected_lang,
                        "tokens": self._count_tokens(chunk),
                        "text": chunk,
                    }
                )
            page_num += 1
        data = {"chunks": chunks_data, "model": self.model, "source": "xlsx"}

        if output_json_path:
            try:
                with open(output_json_path, "w", encoding="utf-8") as json_file:
                    json.dump(data, json_file, indent=4, ensure_ascii=False)
                logging.info("Extracted text from XLSX saved to %s", output_json_path)
            except Exception as e:
                logging.error(
                    "Failed to write JSON output to '%s': %s", output_json_path, e
                )
                raise e
        return data

    # ==============================
    # Generic File Extraction Method
    # ==============================

    def extract_text_from_file(
        self, file_path: str, output_json_path: Optional[str] = None
    ) -> Dict:
        """
        Determines the file type based on its extension and routes extraction
        to the appropriate method.

        Supported file types: PDF, image (png, jpg, jpeg, tiff), DOC, DOCX, PPT, PPTX, CSV, XLSX.

        Args:
            file_path (str): Path to the file.
            output_json_path (Optional[str]): Path to save the JSON output.

        Returns:
            Dict: A dictionary containing the extracted text chunks and model info.
        """
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return self.extract_text_from_pdf(file_path, output_json_path)
        elif ext in [".png", ".jpg", ".jpeg", ".tiff"]:
            return self.extract_text_from_image(file_path, output_json_path)
        elif ext == ".docx":
            return self.extract_text_from_docx(file_path, output_json_path)
        elif ext == ".pptx":
            return self.extract_text_from_pptx(file_path, output_json_path)
        elif ext == ".csv":
            return self.extract_text_from_csv(file_path, output_json_path)
        elif ext == ".xlsx":
            return self.extract_text_from_xlsx(file_path, output_json_path)
        else:
            logging.error("Unsupported file format: %s", ext)
            raise ValueError(f"Unsupported file format: {ext}")


if __name__ == "__main__":
    extractor = DocumentTextEmbeddingsExtractor(model="openai")

    # Example usage for PDF extraction:
    pdf_file = "example.pdf"  # Replace with your PDF file path
    pdf_output = "output_pdf_chunks.json"
    extractor.extract_text_from_pdf(pdf_file, pdf_output)

    # Example usage for Image extraction:
    image_file = "sample_image.png"  # Replace with your image file path
    image_output = "output_image_chunks.json"
    extractor.extract_text_from_image(image_file, image_output)

    # Example usage for DOCX extraction:
    docx_file = "example.docx"  # Replace with your DOCX file path
    docx_output = "output_docx_chunks.json"
    extractor.extract_text_from_docx(docx_file, docx_output)

    # Example usage for DOC extraction (requires textract):
    # doc_file = "example.doc"
    # doc_output = "output_doc_chunks.json"
    # extractor.extract_text_from_doc(doc_file, doc_output)

    # Example usage for PPTX extraction:
    pptx_file = "example.pptx"  # Replace with your PPTX file path
    pptx_output = "output_pptx_chunks.json"
    extractor.extract_text_from_pptx(pptx_file, pptx_output)

    # Example usage for PPT extraction (requires textract):
    # ppt_file = "example.ppt"
    # ppt_output = "output_ppt_chunks.json"
    # extractor.extract_text_from_ppt(ppt_file, ppt_output)

    # Example usage for CSV extraction:
    csv_file = "example.csv"  # Replace with your CSV file path
    csv_output = "output_csv_chunks.json"
    extractor.extract_text_from_csv(csv_file, csv_output)

    # Example usage for XLSX extraction:
    xlsx_file = "example.xlsx"  # Replace with your XLSX file path
    xlsx_output = "output_xlsx_chunks.json"
    extractor.extract_text_from_xlsx(xlsx_file, xlsx_output)

    # Alternatively, to let the extractor auto-detect the file type:
    # file_path = "example.docx"  # Replace with any supported file type
    # output_json = "output_file_chunks.json"
    # extractor.extract_text_from_file(file_path, output_json)
